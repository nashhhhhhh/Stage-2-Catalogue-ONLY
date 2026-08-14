import json
import math
from io import BytesIO
from pathlib import Path

from PIL import Image, ImageFilter, ImageOps
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml import parse_xml


BASE_DIR = Path(__file__).resolve().parent
PPTX_PATH = BASE_DIR / "layout_sources" / "Stage 2 PPT Layout.pptx"
DATA_DIR = BASE_DIR / "frontend" / "static" / "data"
IMAGE_DIR = BASE_DIR / "frontend" / "static" / "images"
CONFIG_PATH = DATA_DIR / "layout_source_config.json"
ASSET_COUNTS_PATH = DATA_DIR / "catalogue_asset_counts.json"

NS = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
SCHEME_COLORS = {"accent4": "8064A2"}
RISK_BY_PREFIX = {"M": "medium", "H": "high", "L": "low", "O": "office"}

DEFAULT_LAYOUTS = {
    "factory": {
        "label": "Factory Layout",
        "slide": 3,
        "picture": "Picture 4",
        "image": "stage2_layout.png",
        "target_width": 7680,
    },
    "incoming-office-level-1": {
        "label": "Incoming Warehouse Office - Level 1",
        "slide": 4,
        "picture": "Picture 10",
        "image": "incoming_warehouse_office_level1.png",
        "target_width": 3600,
    },
    "incoming-office-level-2": {
        "label": "Incoming Warehouse Office - Level 2",
        "slide": 5,
        "picture": "Picture 9",
        "image": "incoming_warehouse_office_level2.png",
        "target_width": 3600,
    },
}
DEFAULT_CONFIG = {
    "source_pptx": str(PPTX_PATH.relative_to(BASE_DIR)).replace("\\", "/"),
    "layouts": DEFAULT_LAYOUTS,
}


def load_layout_source_config():
    config = json.loads(json.dumps(DEFAULT_CONFIG))
    if CONFIG_PATH.exists():
        try:
            saved = json.loads(CONFIG_PATH.read_text(encoding="utf-8"))
        except (json.JSONDecodeError, OSError):
            saved = {}
        config["source_pptx"] = saved.get("source_pptx") or config["source_pptx"]
        saved_layouts = saved.get("layouts") or {}
        for key, defaults in DEFAULT_LAYOUTS.items():
            merged = {**defaults, **(saved_layouts.get(key) or {})}
            merged["slide"] = int(merged.get("slide") or defaults["slide"])
            merged["target_width"] = int(merged.get("target_width") or defaults["target_width"])
            merged["picture"] = str(merged.get("picture") or defaults["picture"]).strip()
            merged["image"] = defaults["image"]
            merged["label"] = str(merged.get("label") or defaults["label"]).strip()
            config["layouts"][key] = merged
    return config


def save_layout_source_config(config):
    DATA_DIR.mkdir(parents=True, exist_ok=True)
    CONFIG_PATH.write_text(json.dumps(config, indent=2), encoding="utf-8")


def load_asset_counts():
    if not ASSET_COUNTS_PATH.exists():
        return {}
    try:
        payload = json.loads(ASSET_COUNTS_PATH.read_text(encoding="utf-8"))
    except (json.JSONDecodeError, OSError):
        return {}
    return {
        code: int(details.get("asset_count") or 0)
        for code, details in payload.items()
        if isinstance(details, dict)
    }


def apply_asset_counts(rooms, asset_counts):
    if not asset_counts:
        return rooms
    for code, room in rooms.items():
        room["assetCount"] = asset_counts.get(code, room.get("assetCount", 0))
    return rooms


def rounded(value):
    return round(value, 4)


def parse_room_shape(shape_name):
    if "_" not in shape_name:
        return None

    code, raw_name = shape_name.split("_", 1)
    code = code.strip()
    risk = RISK_BY_PREFIX.get(code[:1].upper())
    if not code or not risk:
        return None

    room_name = raw_name.replace("_", " ").replace(" . ", " ").strip()
    room_name = " ".join(room_name.split())
    return code, room_name, risk


def find_shape(slide, shape_name):
    for shape in slide.shapes:
        if shape.name == shape_name:
            return shape
    available = ", ".join(shape.name for shape in slide.shapes if "Picture" in shape.name) or "no picture shapes"
    raise ValueError(f"Picture shape '{shape_name}' not found. Available pictures: {available}")


def crop_picture_image(picture):
    image = Image.open(BytesIO(picture.image.blob)).convert("RGBA")
    width, height = image.size
    box = (
        round(picture.crop_left * width),
        round(picture.crop_top * height),
        round((1 - picture.crop_right) * width),
        round((1 - picture.crop_bottom) * height),
    )
    return image.crop(box)


def get_shape_transform(shape):
    transform = shape.element.find(".//a:xfrm", namespaces=NS)
    if transform is None:
        return 0, False, False
    rotation = int(transform.get("rot") or 0) / 60000
    flip_h = str(transform.get("flipH") or "").lower() in {"1", "true"}
    flip_v = str(transform.get("flipV") or "").lower() in {"1", "true"}
    return rotation, flip_h, flip_v


def transform_shape_point(shape, local_x, local_y):
    rotation, flip_h, flip_v = get_shape_transform(shape)
    if flip_h:
        local_x = 1 - local_x
    if flip_v:
        local_y = 1 - local_y

    center_x = shape.left + shape.width / 2
    center_y = shape.top + shape.height / 2
    offset_x = (local_x - 0.5) * shape.width
    offset_y = (local_y - 0.5) * shape.height
    radians = math.radians(rotation)
    return (
        center_x + offset_x * math.cos(radians) - offset_y * math.sin(radians),
        center_y + offset_x * math.sin(radians) + offset_y * math.cos(radians),
    )


def get_rendered_shape_bounds(shape):
    corners = [
        transform_shape_point(shape, local_x, local_y)
        for local_x, local_y in ((0, 0), (1, 0), (1, 1), (0, 1))
    ]
    xs = [point[0] for point in corners]
    ys = [point[1] for point in corners]
    return min(xs), min(ys), max(xs), max(ys)


def export_floorplan(picture, output_path, target_width):
    image = crop_picture_image(picture)
    rotation, flip_h, flip_v = get_shape_transform(picture)
    left, top, right, bottom = get_rendered_shape_bounds(picture)
    rendered_width = max(1, right - left)
    rendered_height = max(1, bottom - top)
    scale = target_width / rendered_width
    unrotated_size = (
        max(1, round(picture.width * scale)),
        max(1, round(picture.height * scale)),
    )
    image = image.resize(unrotated_size, Image.Resampling.LANCZOS)
    if flip_h:
        image = ImageOps.mirror(image)
    if flip_v:
        image = ImageOps.flip(image)
    if rotation % 360:
        image = image.rotate(-rotation, expand=True, resample=Image.Resampling.BICUBIC)

    target_height = max(1, round(target_width * rendered_height / rendered_width))
    if image.size != (target_width, target_height):
        image = image.resize((target_width, target_height), Image.Resampling.LANCZOS)
    image = image.filter(ImageFilter.UnsharpMask(radius=1.2, percent=135, threshold=2))
    output_path.parent.mkdir(parents=True, exist_ok=True)
    image.save(output_path, "PNG", optimize=True)
    return target_width, target_height


def shape_point_to_picture(point, path_width, path_height, shape, picture):
    local_x = int(point.get("x")) / path_width
    local_y = int(point.get("y")) / path_height
    slide_x, slide_y = transform_shape_point(shape, local_x, local_y)
    picture_left, picture_top, picture_right, picture_bottom = get_rendered_shape_bounds(picture)
    return (
        rounded((slide_x - picture_left) / (picture_right - picture_left) * 100),
        rounded((slide_y - picture_top) / (picture_bottom - picture_top) * 100),
    )


def get_shape_rotation(shape):
    return get_shape_transform(shape)[0]


def extract_rotated_rect_path(shape, picture):
    rotation = get_shape_rotation(shape)
    if rotation % 360 == 0:
        return None

    picture_left, picture_top, picture_right, picture_bottom = get_rendered_shape_bounds(picture)
    picture_width = picture_right - picture_left
    picture_height = picture_bottom - picture_top
    corners = [
        (
            rounded((slide_x - picture_left) / picture_width * 100),
            rounded((slide_y - picture_top) / picture_height * 100),
        )
        for slide_x, slide_y in (
            transform_shape_point(shape, local_x, local_y)
            for local_x, local_y in ((0, 0), (1, 0), (1, 1), (0, 1))
        )
    ]

    return " ".join(
        [f"M {corners[0][0]} {corners[0][1]}"]
        + [f"L {x} {y}" for x, y in corners[1:]]
        + ["Z"]
    )


def extract_svg_path(shape, picture):
    path = shape.element.find(".//a:custGeom/a:pathLst/a:path", namespaces=NS)
    if path is None:
        return extract_rotated_rect_path(shape, picture)

    path_width = int(path.get("w") or shape.width)
    path_height = int(path.get("h") or shape.height)
    commands = []

    for command in path:
        tag = command.tag.rsplit("}", 1)[-1]
        points = command.findall("a:pt", namespaces=NS)
        coords = [
            shape_point_to_picture(point, path_width, path_height, shape, picture)
            for point in points
        ]

        if tag == "moveTo" and coords:
            commands.append(f"M {coords[0][0]} {coords[0][1]}")
        elif tag == "lnTo" and coords:
            commands.append(f"L {coords[0][0]} {coords[0][1]}")
        elif tag == "cubicBezTo" and len(coords) == 3:
            commands.append("C " + " ".join(f"{x} {y}" for x, y in coords))
        elif tag == "quadBezTo" and len(coords) == 2:
            commands.append("Q " + " ".join(f"{x} {y}" for x, y in coords))
        elif tag == "close":
            commands.append("Z")

    return " ".join(commands) or None


def get_theme_colors(slide):
    colors = dict(SCHEME_COLORS)
    try:
        master = slide.slide_layout.slide_master
        theme_part = master.part.part_related_by(RT.THEME)
        theme = parse_xml(theme_part.blob)
        color_scheme = theme.find(".//a:themeElements/a:clrScheme", namespaces=NS)
        if color_scheme is not None:
            for scheme_entry in color_scheme:
                if len(scheme_entry) == 0:
                    continue
                color_node = scheme_entry[0]
                tag = color_node.tag.rsplit("}", 1)[-1]
                value = color_node.get("val") if tag == "srgbClr" else color_node.get("lastClr")
                if value:
                    colors[scheme_entry.tag.rsplit("}", 1)[-1]] = value.upper()

        color_map = master.element.find(".//p:clrMap", namespaces={
            **NS,
            "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
        })
        if color_map is not None:
            for alias, target in color_map.attrib.items():
                if target in colors:
                    colors[alias] = colors[target]
    except (AttributeError, KeyError, ValueError):
        pass
    return colors


def resolve_color(container, default, theme_colors=None):
    if container is None or len(container) == 0:
        return default
    color = container[0]
    tag = color.tag.rsplit("}", 1)[-1]
    if tag == "srgbClr":
        return color.get("val") or default
    if tag == "schemeClr":
        return (theme_colors or SCHEME_COLORS).get(color.get("val"), default)
    if tag == "sysClr":
        return color.get("lastClr") or default
    return default


def extract_style(shape, risk, theme_colors=None):
    shape_properties = shape.element.find("p:spPr", namespaces={
        **NS,
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    })
    solid_fill = shape_properties.find("a:solidFill", namespaces=NS) if shape_properties is not None else None
    line = shape_properties.find("a:ln", namespaces=NS) if shape_properties is not None else None
    line_fill = line.find("a:solidFill", namespaces=NS) if line is not None else None
    color_defaults = {
        "medium": "0EA5E9",
        "high": "22C55E",
        "low": "EAB308",
        "office": "8064A2",
    }
    fill_color = resolve_color(solid_fill, color_defaults[risk], theme_colors)
    stroke_color = resolve_color(line_fill, fill_color, theme_colors)
    color_node = solid_fill[0] if solid_fill is not None and len(solid_fill) else None
    alpha = color_node.find("a:alpha", namespaces=NS) if color_node is not None else None
    fill_opacity = int(alpha.get("val")) / 100000 if alpha is not None else 1
    line_width = int(line.get("w") or 12700) / 12700 if line is not None else 1
    return {
        "fillColor": f"#{fill_color}",
        "fillOpacity": rounded(fill_opacity),
        "strokeColor": f"#{stroke_color}",
        "strokeWidth": rounded(line_width),
    }


def extract_room(shape, picture, level=None, theme_colors=None):
    parsed = parse_room_shape(shape.name)
    if not parsed:
        return None

    code, room_name, risk = parsed
    shape_left, shape_top, shape_right, shape_bottom = get_rendered_shape_bounds(shape)
    picture_left, picture_top, picture_right, picture_bottom = get_rendered_shape_bounds(picture)
    picture_width = picture_right - picture_left
    picture_height = picture_bottom - picture_top
    room = {
        "code": code,
        "name": room_name,
        "risk": risk,
        "interactive": True,
        "left": rounded((shape_left - picture_left) / picture_width * 100),
        "top": rounded((shape_top - picture_top) / picture_height * 100),
        "width": rounded((shape_right - shape_left) / picture_width * 100),
        "height": rounded((shape_bottom - shape_top) / picture_height * 100),
        "svgPath": extract_svg_path(shape, picture),
        **extract_style(shape, risk, theme_colors),
    }
    if level:
        room["level"] = level
    return room


def extract_layout(presentation, key, config):
    slide = presentation.slides[config["slide"] - 1]
    picture = find_shape(slide, config["picture"])
    theme_colors = get_theme_colors(slide)
    image_size = export_floorplan(
        picture,
        IMAGE_DIR / config["image"],
        config["target_width"],
    )
    level = None
    if key == "incoming-office-level-1":
        level = "Incoming Warehouse Office - Level 1"
    elif key == "incoming-office-level-2":
        level = "Incoming Warehouse Office - Level 2"

    rooms = {}
    for shape in slide.shapes:
        if shape.shape_type not in {
            MSO_SHAPE_TYPE.AUTO_SHAPE,
            MSO_SHAPE_TYPE.FREEFORM,
        }:
            continue
        room = extract_room(shape, picture, level, theme_colors)
        if room:
            rooms[room["code"]] = room

    return dict(sorted(rooms.items())), {
        "label": config.get("label", key),
        "source_pptx": config.get("source_pptx", str(PPTX_PATH.relative_to(BASE_DIR)).replace("\\", "/")),
        "slide": config["slide"],
        "picture": config["picture"],
        "image": config["image"],
        "target_width": config["target_width"],
        "width": image_size[0],
        "height": image_size[1],
        "aspect": rounded(image_size[0] / image_size[1]),
    }


def main():
    config = load_layout_source_config()
    pptx_path = BASE_DIR / config["source_pptx"]
    if not pptx_path.exists():
        raise FileNotFoundError(f"PPTX file not found: {pptx_path}")

    presentation = Presentation(pptx_path)
    extracted = {}
    metadata = {}
    for key, layout_config in config["layouts"].items():
        metadata_config = {**layout_config, "source_pptx": config["source_pptx"]}
        extracted[key], metadata[key] = extract_layout(presentation, key, metadata_config)
        if not extracted[key]:
            label = layout_config.get("label", key)
            raise ValueError(
                f"No named room shapes found for {label} on slide {layout_config['slide']}. "
                "Name the overlay shapes like O-01_Room_Name or L-01_Room_Name, then rebuild."
            )

    asset_counts = load_asset_counts()
    apply_asset_counts(extracted["factory"], asset_counts)
    apply_asset_counts(extracted["incoming-office-level-1"], asset_counts)
    apply_asset_counts(extracted["incoming-office-level-2"], asset_counts)

    DATA_DIR.mkdir(parents=True, exist_ok=True)
    save_layout_source_config(config)
    (DATA_DIR / "room_shapes.json").write_text(
        json.dumps(extracted["factory"], indent=2),
        encoding="utf-8",
    )
    (DATA_DIR / "office_layout_shapes.json").write_text(
        json.dumps({
            "incoming-office-level-1": extracted["incoming-office-level-1"],
            "incoming-office-level-2": extracted["incoming-office-level-2"],
        }, indent=2),
        encoding="utf-8",
    )
    (DATA_DIR / "layout_metadata.json").write_text(
        json.dumps(metadata, indent=2),
        encoding="utf-8",
    )

    print(
        "Rebuilt layouts:",
        f"{len(extracted['factory'])} factory rooms,",
        f"{len(extracted['incoming-office-level-1'])} level-1 rooms,",
        f"{len(extracted['incoming-office-level-2'])} level-2 rooms.",
    )


if __name__ == "__main__":
    main()
