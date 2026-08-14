import io
import json
import sys
import unittest
from pathlib import Path
from unittest.mock import mock_open, patch

from PIL import Image
import fitz
from pptx import Presentation
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches


ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

import app as app_module
import custom_catalogue_service as custom_service
import machine_catalogue_service as machine_service


class CatalogueAppTests(unittest.TestCase):
    def setUp(self):
        self.refresh_patch = patch.object(app_module, "schedule_catalogue_refresh")
        self.custom_refresh_patch = patch.object(app_module, "schedule_custom_catalogue_refresh")
        self.refresh_patch.start()
        self.custom_refresh_patch.start()
        self.client = app_module.app.test_client()

    def tearDown(self):
        self.custom_refresh_patch.stop()
        self.refresh_patch.stop()

    def test_main_pages_and_read_apis(self):
        expected_statuses = {
            "/": 200,
            "/layout": 200,
            "/catalogue/M-01": 200,
            "/catalogue/manage": 200,
            "/catalogue/manage/machines": 200,
            "/catalogue/manage/rooms": 200,
            "/catalogue/manage/create": 200,
            "/api/catalogue/layout-source": 200,
            "/api/catalogue/management/catalogues": 200,
            "/api/catalogue/rooms": 200,
            "/api/catalogue/room-catalogues": 200,
            "/api/catalogue/machines": 200,
            "/api/catalogue/custom": 200,
            "/api/catalogue/status/not-a-risk": 400,
        }
        for route, expected_status in expected_statuses.items():
            with self.subTest(route=route):
                response = self.client.get(route)
                self.assertEqual(response.status_code, expected_status)
                response.close()

    def test_fast_pdf_link_parser_keeps_expected_room_pages(self):
        expected_pages = {
            "medium": [13, 15, 22, 24, 25, 36, 37, 39, 41, 42],
            "high": [14, 15, 18, 19, 20, 21, 22, 23, 24, 25, 26, 27, 28, 29, 30, 31, 32, 33, 38],
            "low": [17, 22, 23, 24, 25, 26, 27, 30, 34, 38, 41, 43, 45, 46, 50, 51, 53, 55, 56, 58, 65, 67, 69, 70, 71, 72, 74, 76, 77, 78, 80, 83, 84, 86, 88, 90, 92, 94, 95, 97, 98, 99, 102, 103, 104, 105, 106, 107, 109, 112, 114, 116, 119],
            "office": [15, 17, 19, 20, 21, 22, 23, 24, 25, 26, 28, 33, 36, 38, 39],
        }
        for risk_area, expected in expected_pages.items():
            with self.subTest(risk_area=risk_area):
                pdf_url = f"/static/catalogue/current_{risk_area}_catalogue.pdf"
                self.assertEqual(app_module.get_pdf_contents_link_pages(pdf_url), expected)

    def test_built_in_mapping_uses_headers_even_when_contents_is_incomplete(self):
        catalogue_map = {
            "L-49": {"name": "LR War RM", "page": 107, "pdf": "/static/catalogue/current_low_catalogue.pdf"},
            "L-50": {"name": "Staff Access LR Area-2", "page": 108, "pdf": "/static/catalogue/current_low_catalogue.pdf"},
            "L-51": {"name": "LR Area", "page": 111, "pdf": "/static/catalogue/current_low_catalogue.pdf"},
        }
        link_rooms = {
            **catalogue_map,
            "L-50": {**catalogue_map["L-50"], "page": 109},
        }
        header_rooms = {
            **link_rooms,
            "L-51": {**catalogue_map["L-51"], "page": 112},
        }
        app_module._CATALOGUE_ROOM_PAGE_CACHE.clear()
        with (
            patch.object(app_module, "get_catalogue_pdf_path", return_value="test.pdf"),
            patch.object(app_module, "get_file_signature", return_value=(1, 1)),
            patch("builtins.open", mock_open(read_data=b"%PDF-test")),
            patch.object(app_module, "get_pdf_contents_links", return_value=[{"text": "L-50", "page": 109}]),
            patch.object(app_module, "match_pdf_pages_to_rooms", return_value=(link_rooms, 1)),
            patch.object(app_module, "match_pdf_headers_to_rooms", return_value=(header_rooms, 3)),
        ):
            mapped = app_module.apply_catalogue_toc_pages(catalogue_map)
        self.assertEqual(mapped["L-49"]["page"], 107)
        self.assertEqual(mapped["L-50"]["page"], 109)
        self.assertEqual(mapped["L-51"]["page"], 112)

    def test_live_asset_counts_match_checked_in_fallback(self):
        fallback_path = ROOT / "frontend" / "static" / "data" / "catalogue_asset_counts.json"
        fallback = json.loads(fallback_path.read_text(encoding="utf-8"))["counts"]
        live = app_module.build_catalogue_asset_counts()
        self.assertEqual(set(live), set(fallback))
        for room_code, details in live.items():
            with self.subTest(room_code=room_code):
                self.assertEqual(details["asset_count"], fallback[room_code]["asset_count"])
                self.assertEqual(details["source"], "catalogue_pdf")

    def test_custom_powerpoint_accepts_friendly_shape_names(self):
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        image_stream = io.BytesIO()
        Image.new("RGB", (640, 360), "white").save(image_stream, format="PNG")
        image_stream.seek(0)
        picture = slide.shapes.add_picture(
            image_stream,
            Inches(0),
            Inches(0),
            width=Inches(10),
            height=Inches(5.625),
        )
        names = ["Canteen Area", "Dish Store & Washing Area", "Kitchen_Area"]
        for index, name in enumerate(names):
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(1 + index * 2),
                Inches(1),
                Inches(1.5),
                Inches(1.5),
            )
            shape.name = name

        rooms = custom_service.extract_rooms_from_slide(slide, picture)
        self.assertEqual(list(rooms), ["C-01", "C-02", "C-03"])
        self.assertEqual(
            [room["name"] for room in rooms.values()],
            ["Canteen Area", "Dish Store & Washing Area", "Kitchen Area"],
        )

    def test_custom_powerpoint_accepts_multiple_map_slides(self):
        presentation = Presentation()
        extracted = {}
        next_index = 1
        for slide_number in (1, 2):
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            image_stream = io.BytesIO()
            Image.new("RGB", (640, 360), "white").save(image_stream, format="PNG")
            image_stream.seek(0)
            picture = slide.shapes.add_picture(
                image_stream,
                Inches(0),
                Inches(0),
                width=Inches(10),
                height=Inches(5.625),
            )
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(slide_number),
                Inches(1),
                Inches(1.5),
                Inches(1.5),
            )
            shape.name = f"Map Section {slide_number}"
            slide_rooms = custom_service.extract_rooms_from_slide(slide, picture, next_index)
            for code, room in slide_rooms.items():
                extracted[code] = {**room, "layoutKey": f"slide-{slide_number}"}
            next_index += len(slide_rooms)

        self.assertEqual(custom_service.parse_slide_numbers("all", 2), [1, 2])
        self.assertEqual(custom_service.parse_slide_numbers("2, 1-2", 2), [2, 1])
        self.assertEqual(list(extracted), ["C-01", "C-02"])
        self.assertEqual(extracted["C-01"]["layoutKey"], "slide-1")
        self.assertEqual(extracted["C-02"]["layoutKey"], "slide-2")

    def test_custom_powerpoint_combines_adjacent_map_picture_tiles(self):
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        for tile_index in range(2):
            image_stream = io.BytesIO()
            Image.new("RGB", (320, 360), "white").save(image_stream, format="PNG")
            image_stream.seek(0)
            slide.shapes.add_picture(
                image_stream,
                Inches(tile_index * 5),
                Inches(0),
                width=Inches(5),
                height=Inches(5.625),
            )
            room = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(1 + tile_index * 5),
                Inches(1),
                Inches(2),
                Inches(2),
            )
            room.name = f"U-0{tile_index + 1}_Tile_{tile_index + 1}"

        pictures = custom_service.find_map_pictures(slide)
        rooms = custom_service.extract_rooms_from_slide(slide, pictures)
        self.assertEqual(len(pictures), 2)
        self.assertEqual(list(rooms), ["U-01", "U-02"])
        self.assertGreaterEqual(min(room["left"] for room in rooms.values()), 0)
        self.assertLessEqual(max(room["left"] + room["width"] for room in rooms.values()), 100)

        with patch.object(Image.Image, "save"):
            exported_size = custom_service.export_floorplan(
                pictures,
                ROOT / "tiled-layout-test.png",
                1000,
            )
        self.assertEqual(exported_size, (1000, 562))

    def test_custom_catalogue_creation_exports_every_selected_map_slide(self):
        presentation = Presentation()
        for slide_number in (1, 2):
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            image_stream = io.BytesIO()
            Image.new("RGB", (320, 180), "white").save(image_stream, format="PNG")
            image_stream.seek(0)
            slide.shapes.add_picture(
                image_stream,
                Inches(0),
                Inches(0),
                width=Inches(10),
                height=Inches(5.625),
            )
            room = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(slide_number),
                Inches(1),
                Inches(2),
                Inches(1.5),
            )
            room.name = f"Section {slide_number} Room"

        with (
            patch.object(custom_service, "Presentation", return_value=presentation),
            patch.object(custom_service.Path, "mkdir"),
            patch.object(custom_service.Path, "open", mock_open()),
            patch.object(custom_service.shutil, "copyfile"),
            patch.object(custom_service, "export_floorplan", return_value=(800, 450)) as export,
            patch.object(custom_service, "load_custom_catalogues", return_value=[]),
            patch.object(custom_service, "save_custom_catalogues"),
        ):
            catalogue = custom_service.create_custom_catalogue(
                "Two Map Slides",
                "",
                ROOT / "two-map-slides.pptx",
                b"",
                {"slug": "two-map-slides", "slide_numbers": "all", "target_width": 800},
            )

        self.assertEqual(catalogue["slide_numbers"], [1, 2])
        self.assertEqual([layout["key"] for layout in catalogue["layouts"]], ["slide-1", "slide-2"])
        self.assertEqual(catalogue["rooms"]["C-01"]["layoutKey"], "slide-1")
        self.assertEqual(catalogue["rooms"]["C-02"]["layoutKey"], "slide-2")
        self.assertEqual(export.call_count, 2)
        self.assertEqual(export.call_args_list[0].args[1].name, "layout.png")
        self.assertEqual(export.call_args_list[1].args[1].name, "layout-slide-2.png")

    def test_custom_powerpoint_preserves_rotation_and_theme_color(self):
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        image_stream = io.BytesIO()
        Image.new("RGB", (80, 120), "white").save(image_stream, format="PNG")
        image_stream.seek(0)
        picture = slide.shapes.add_picture(
            image_stream,
            Inches(2),
            Inches(1),
            width=Inches(2),
            height=Inches(3),
        )
        picture.rotation = 90

        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(2.5),
            Inches(1.5),
            Inches(1),
            Inches(2),
        )
        shape.name = "C-01_Rotated Room"
        shape.rotation = 90
        shape.fill.solid()
        shape.fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_3

        room = custom_service.extract_rooms_from_slide(slide, picture)["C-01"]
        theme_colors = custom_service.get_theme_colors(slide)
        self.assertAlmostEqual(room["left"], 16.6667, places=3)
        self.assertAlmostEqual(room["top"], 25, places=3)
        self.assertAlmostEqual(room["width"], 66.6667, places=3)
        self.assertAlmostEqual(room["height"], 50, places=3)
        self.assertIsNotNone(room["svgPath"])
        self.assertEqual(room["fillColor"], f"#{theme_colors['accent3']}")

        output_path = ROOT / "rotated-layout-test.png"
        with patch.object(Image.Image, "save") as save_image:
            exported_size = custom_service.export_floorplan(picture, output_path, 600)
            self.assertEqual(exported_size, (600, 400))
        save_image.assert_called_once_with(output_path, "PNG", optimize=True)

    def test_custom_room_names_match_google_doc_contents_links(self):
        rooms = {
            "C-01": {"code": "C-01", "name": "Canteen Area"},
            "C-02": {"code": "C-02", "name": "Dish Store & Washing Area"},
            "C-03": {"code": "C-03", "name": "Kitchen Area"},
        }
        links = [
            {"page": 4, "text": "1 canteen area 4"},
            {"page": 7, "text": "2 dish store washing area 7"},
            {"page": 10, "text": "3 kitchen area 10"},
        ]
        matched, match_count = custom_service.match_pdf_pages_to_rooms(rooms, links)
        self.assertEqual(match_count, 3)
        self.assertEqual([matched[code]["page"] for code in rooms], [4, 7, 10])

    def test_added_catalogue_room_overrides_are_scoped_and_reusable(self):
        rooms = {
            "C-01": {"code": "C-01", "name": "Original Room", "page": 4},
        }
        updated = custom_service.apply_custom_room_overrides(rooms, {
            "C-01": {"name": "Renamed Room", "page": 9},
        })
        self.assertEqual(updated["C-01"]["name"], "Renamed Room")
        self.assertEqual(updated["C-01"]["page"], 9)
        self.assertEqual(updated["C-01"]["defaultName"], "Original Room")
        self.assertEqual(updated["C-01"]["defaultPage"], 4)

    def test_custom_room_pages_are_validated_against_pdf_title_headers(self):
        pdf = fitz.open()
        contents = pdf.new_page()
        contents.insert_text((72, 72), "Contents\nMaintenance Office ........ 2")
        room_page = pdf.new_page()
        room_page.insert_text((72, 72), "Maintenance Office + Laundry Rooms:\n1. Maintenance Office")
        pdf_bytes = pdf.tobytes()
        pdf.close()
        rooms = {"C-01": {"code": "C-01", "name": "Maintenance Office", "page": 1}}
        matched, match_count = custom_service.match_pdf_headers_to_rooms(rooms, pdf_bytes)
        self.assertEqual(match_count, 1)
        self.assertEqual(matched["C-01"]["page"], 2)

    def test_custom_catalogue_summary_exposes_layout_for_main_selector(self):
        item = {
            "slug": "test-layout",
            "title": "Test Layout",
            "rooms": {"C-01": {}},
            "layout": {
                "image": "/static/custom_catalogues/test-layout/layout.png",
                "width": 1200,
                "height": 600,
                "aspect": 2,
            },
        }
        with patch.object(app_module, "load_custom_catalogues", return_value=[item]):
            response = self.client.get("/api/catalogue/custom")
        self.assertEqual(response.status_code, 200)
        payload = response.get_json()[0]
        self.assertEqual(payload["slug"], "test-layout")
        self.assertEqual(payload["layout"]["aspect"], 2)
        self.assertEqual(payload["view_url"], "/catalogue/custom/test-layout")

    def test_custom_catalogue_assets_disable_stale_browser_cache(self):
        catalogue = {
            "slug": "test-layout",
            "title": "Test Layout",
            "layout": {
                "image": "/static/custom_catalogues/test-layout/layout.png",
                "width": 1200,
                "height": 500,
                "aspect": 2.4,
            },
            "rooms": {},
        }
        with patch.object(app_module, "load_custom_catalogues", return_value=[catalogue]):
            api_response = self.client.get("/api/catalogue/custom")
        self.assertIn("no-store", api_response.headers["Cache-Control"])

        with patch.object(app_module, "get_custom_catalogue", return_value=catalogue):
            item_response = self.client.get("/api/catalogue/custom/test-layout")
        self.assertIn("no-store", item_response.headers["Cache-Control"])

        layout_response = self.client.get("/layout")
        self.assertIn("no-store", layout_response.headers["Cache-Control"])
        layout_response.close()

        layout_page = (ROOT / "frontend" / "layout_map.html").read_text(encoding="utf-8")
        custom_page = (ROOT / "frontend" / "custom_catalogue_view.html").read_text(encoding="utf-8")
        self.assertIn("imageVersion", layout_page)
        self.assertIn("imageVersion", custom_page)
        self.assertIn("mapLayouts", custom_page)
        self.assertIn("room.layoutKey", custom_page)

    def test_management_catalogue_filter_includes_added_and_future_catalogues(self):
        added = {
            "slug": "stage-2-canteen",
            "title": "Stage 2 Canteen",
            "doc_url": "https://docs.google.com/document/d/canteen/edit",
            "source_pptx": "data/custom_catalogue_uploads/stage-2-canteen.pptx",
            "slide_number": 2,
            "picture_name": "Picture 1",
            "room_count": 3,
            "layout": {"width": 3600, "height": 1800, "aspect": 2},
        }
        with (
            patch.object(app_module, "load_custom_catalogues", return_value=[added]),
            patch.object(app_module, "get_pptx_slide_count", return_value=5),
        ):
            response = self.client.get("/api/catalogue/management/catalogues")
        payload = response.get_json()["catalogues"]
        self.assertEqual(response.status_code, 200)
        self.assertEqual(len(payload), 5)
        canteen = next(item for item in payload if item["id"] == "custom:stage-2-canteen")
        self.assertEqual(canteen["title"], "Stage 2 Canteen")
        self.assertEqual(canteen["layouts"]["custom"]["room_count"], 3)

    def test_custom_management_update_replaces_same_catalogue_and_preserves_overrides(self):
        current = {
            "slug": "stage-2-canteen",
            "title": "Stage 2 Canteen",
            "doc_url": "https://docs.google.com/document/d/canteen/edit",
            "source_pptx": "data/custom_catalogue_uploads/stage-2-canteen.pptx",
            "created_at": "2026-08-04 10:00:00",
            "slide_number": 1,
            "picture_name": "Picture 1",
            "room_overrides": {"C-01": {"name": "Main Canteen"}},
            "layout": {
                "pdf": "/static/custom_catalogues/stage-2-canteen/catalogue.pdf",
                "width": 3600,
            },
        }
        rebuilt = {"slug": "stage-2-canteen", "title": "Stage 2 Canteen", "room_count": 3}
        with (
            patch.object(app_module, "get_custom_catalogue", return_value=current),
            patch.object(app_module.os.path, "isfile", return_value=True),
            patch("builtins.open", mock_open(read_data=b"%PDF-test")),
            patch.object(app_module, "create_custom_catalogue", return_value=rebuilt) as create,
        ):
            response = self.client.post("/api/catalogue/custom/stage-2-canteen/update", data={
                "title": "Stage 2 Canteen",
                "doc_url": current["doc_url"],
                "slide_number": "2",
                "picture_name": "Picture 9",
                "target_width": "4200",
            })
        self.assertEqual(response.status_code, 200)
        options = create.call_args.args[4]
        self.assertEqual(options["slug"], "stage-2-canteen")
        self.assertEqual(options["slide_number"], "2")
        self.assertEqual(options["picture_name"], "Picture 9")
        self.assertEqual(options["room_overrides"], current["room_overrides"])

    def test_management_page_and_layout_modal_share_catalogue_first_editor(self):
        management_response = self.client.get("/catalogue/manage")
        layout_response = self.client.get("/layout")
        management = management_response.get_data(as_text=True)
        layout = layout_response.get_data(as_text=True)
        self.assertIn('id="catalogue-select"', management)
        self.assertIn("Select a catalogue...", management)
        self.assertIn("/api/catalogue/management/catalogues", management)
        self.assertIn('src="/catalogue/manage?embedded=1"', layout)
        management_response.close()
        layout_response.close()

    def test_room_catalogue_filter_includes_added_catalogues_and_scopes_duplicate_codes(self):
        added = {
            "slug": "test-layout",
            "title": "Test Layout",
            "room_count": 1,
            "rooms": {
                "C-01": {"code": "C-01", "name": "Test Room", "page": 7},
            },
        }
        with patch.object(app_module, "load_custom_catalogues", return_value=[added]):
            catalogue_response = self.client.get("/api/catalogue/room-catalogues")
        catalogue_ids = {item["id"] for item in catalogue_response.get_json()}
        self.assertTrue({"area:medium", "area:high", "area:low", "area:office", "custom:test-layout"}.issubset(catalogue_ids))

        with patch.object(app_module, "get_custom_catalogue", return_value=added):
            room_response = self.client.get("/api/catalogue/rooms?catalogue=custom:test-layout")
        rows = room_response.get_json()
        self.assertEqual(len(rows), 1)
        self.assertEqual(rows[0]["catalogue_id"], "custom:test-layout")
        self.assertEqual(rows[0]["code"], "C-01")

    def test_added_catalogue_room_save_uses_catalogue_identity(self):
        saved = {"catalogue_id": "custom:test-layout", "code": "C-01", "name": "Renamed", "page": 8}
        with patch.object(app_module, "upsert_custom_catalogue_room", return_value=saved) as upsert:
            response = self.client.post("/api/catalogue/rooms", json={
                "catalogue_id": "custom:test-layout",
                "code": "C-01",
                "name": "Renamed",
                "page": 8,
            })
        self.assertEqual(response.status_code, 200)
        self.assertEqual(response.get_json()["room"]["name"], "Renamed")
        upsert.assert_called_once()

    def test_custom_catalogue_rebuild_api_reports_updated_room_count(self):
        rebuilt = {
            "slug": "test-layout",
            "title": "Test Layout",
            "room_count": 5,
        }
        with patch.object(app_module, "rebuild_custom_catalogue", return_value=rebuilt):
            response = self.client.post("/api/catalogue/custom/test-layout/rebuild")
        self.assertEqual(response.status_code, 200)
        self.assertEqual(response.get_json()["catalogue"]["room_count"], 5)

    def test_custom_view_uses_pdfjs_canvas_instead_of_native_iframe(self):
        catalogue = {
            "slug": "test-layout",
            "title": "Test Layout",
            "room_count": 1,
            "slide_number": 1,
            "picture_name": "Picture 1",
            "layout": {
                "image": "/static/custom_catalogues/test-layout/layout.png",
                "pdf": "/static/custom_catalogues/test-layout/catalogue.pdf",
                "aspect": 2,
            },
            "rooms": {
                "C-01": {
                    "code": "C-01",
                    "name": "Test Room",
                    "left": 10,
                    "top": 10,
                    "width": 20,
                    "height": 20,
                    "page": 4,
                }
            },
        }
        with patch.object(app_module, "get_custom_catalogue", return_value=catalogue):
            response = self.client.get("/catalogue/custom/test-layout?room=C-01")
        html = response.get_data(as_text=True)
        self.assertEqual(response.status_code, 200)
        self.assertIn('id="pdf-pages"', html)
        self.assertIn("wrapper.className = 'pdf-page'", html)
        self.assertIn('/static/vendor/pdfjs/pdf.min.js', html)
        self.assertIn('id="catalogue-switcher"', html)
        self.assertIn('id="refresh-button"', html)
        self.assertIn('catalogue-help-button', html)
        self.assertIn('/static/catalogue_shell.js', html)
        self.assertIn("searchParams.set('_refresh'", html)
        self.assertIn('hardRefreshToken', html)
        self.assertNotIn('Custom Area', html)
        self.assertNotIn('id="pdf-frame"', html)

    def test_area_catalogue_has_direct_navigation_and_smooth_room_scroll(self):
        response = self.client.get("/catalogue/M-01")
        html = response.get_data(as_text=True)
        self.assertEqual(response.status_code, 200)
        self.assertIn('id="catalogue-switcher"', html)
        self.assertIn('href="/layout"', html)
        self.assertIn('href="/catalogue/manage/rooms"', html)
        self.assertIn('id="refresh-button"', html)
        self.assertIn('catalogue-help-button', html)
        self.assertIn("instant ? 'auto' : 'smooth'", html)
        self.assertIn('/static/catalogue_shell.js', html)
        self.assertIn("searchParams.set('_refresh'", html)
        self.assertIn('hardRefreshToken', html)

    def test_area_refresh_now_uses_the_saved_google_doc_source(self):
        metadata = {
            "current_file": "current_medium_catalogue.pdf",
            "source_url": "https://docs.google.com/document/d/doc-id/edit",
            "last_checked": "2026-08-04 12:00:00",
            "refresh_error": "",
        }
        with (
            patch.object(app_module, "load_catalogue_metadata", return_value=metadata),
            patch.object(app_module, "refresh_catalogue_from_source", return_value=False) as refresh,
            patch.object(app_module, "get_catalogue_pdf_version", return_value="version-1"),
        ):
            response = self.client.post("/api/catalogue/medium/refresh")
        self.assertEqual(response.status_code, 200)
        self.assertFalse(response.get_json()["changed"])
        refresh.assert_called_once_with("medium", force=True)

    def test_catalogue_shell_has_one_unified_list_and_first_visit_guide(self):
        response = self.client.get("/static/catalogue_shell.js")
        script = response.get_data(as_text=True)
        self.assertEqual(response.status_code, 200)
        self.assertIn("Welcome to Catalogue Explorer", script)
        self.assertIn("catalogue-guide-seen-v1", script)
        self.assertNotIn("Custom Catalogues", script)
        response.close()

    def test_layout_cover_uses_catalogue_explorer_shell(self):
        response = self.client.get("/layout")
        html = response.get_data(as_text=True)
        self.assertEqual(response.status_code, 200)
        self.assertIn('class="layout-shell"', html)
        self.assertIn("Interactive floor plans · live equipment catalogues", html)
        self.assertIn('href="/catalogue/manage/machines"', html)
        self.assertIn('href="/catalogue/manage/rooms"', html)
        self.assertIn('href="/catalogue/manage/create"', html)
        self.assertNotIn('/shared/navbar.html', html)
        self.assertGreaterEqual(html.count("fitMode: 'full'"), 2)
        self.assertIn('id="map-zoom-in"', html)
        self.assertIn('id="map-zoom-out"', html)
        self.assertIn('id="map-fit"', html)
        self.assertIn('id="map-lock"', html)
        self.assertIn("catalogue-map-adjustments-v1", html)
        self.assertIn("document.addEventListener('pointermove'", html)
        self.assertIn("mapWrapper.addEventListener('wheel'", html)
        self.assertIn("event.button !== 0 && event.button !== 1", html)
        self.assertIn("saved.locked === true", html)
        self.assertIn("event.ctrlKey", html)
        self.assertNotIn("setPointerCapture", html)
        self.assertIn("room.interactive !== false", html)
        self.assertIn("function getRenderedRoomBounds()", html)
        response.close()

    def test_custom_google_doc_refresh_rebuilds_when_pdf_changes(self):
        catalogue = {
            "slug": "test-layout",
            "doc_url": "https://docs.google.com/document/d/doc-id/edit",
            "layout": {"pdf": "/static/custom_catalogues/test-layout/catalogue.pdf"},
        }
        app_module._CUSTOM_CATALOGUE_REFRESH_STATE.clear()
        with (
            patch.object(app_module, "get_custom_catalogue", return_value=catalogue),
            patch.object(app_module, "download_google_doc_pdf", return_value=(b"%PDF-new", "doc-id")),
            patch.object(app_module.os.path, "exists", return_value=False),
            patch.object(app_module, "rebuild_custom_catalogue") as rebuild,
            patch.object(app_module, "update_custom_catalogue_metadata") as update_metadata,
        ):
            changed = app_module.refresh_custom_catalogue_from_source("test-layout", force=True)
        self.assertTrue(changed)
        rebuild.assert_called_once_with("test-layout", pdf_bytes=b"%PDF-new")
        self.assertEqual(update_metadata.call_args.args[1]["refresh_error"], "")

    def test_machine_and_room_saves_are_round_trip_safe(self):
        machine_payload = machine_service.empty_payload()
        with (
            patch.object(machine_service, "load_machine_catalogue", return_value=machine_payload),
            patch.object(machine_service, "save_machine_catalogue") as save_machines,
        ):
            machine = machine_service.upsert_machine({
                "machine_name": "Test Machine",
                "room_code": "h1",
                "quantity": "2",
            })
            self.assertEqual(machine["room_code"], "H-01")
            save_machines.assert_called_once_with(machine_payload)

        room_payload = machine_service.empty_room_overrides()
        with (
            patch.object(machine_service, "load_room_overrides", return_value=room_payload),
            patch.object(machine_service, "save_room_overrides") as save_rooms,
        ):
            room = machine_service.upsert_room_override({
                "code": "l2",
                "name": "Test Room",
                "page": "7",
            })
            self.assertEqual(room["code"], "L-02")
            self.assertEqual(room["page"], 7)
            save_rooms.assert_called_once_with(room_payload)

    def test_failed_custom_catalogue_upload_removes_incoming_file(self):
        with (
            patch.object(app_module, "DATA_DIR", "test-data"),
            patch.object(app_module.os, "makedirs"),
            patch("werkzeug.datastructures.FileStorage.save"),
            patch.object(app_module.os, "remove") as remove_file,
            patch.object(app_module, "download_google_doc_pdf", return_value=(b"%PDF-test", "doc-id")),
            patch.object(app_module, "create_custom_catalogue", side_effect=ValueError("Invalid test deck")),
        ):
            response = self.client.post(
                "/api/catalogue/custom",
                data={
                    "title": "Test Catalogue",
                    "doc_url": "https://docs.google.com/document/d/doc-id/edit",
                    "pptx_file": (io.BytesIO(b"not-a-real-pptx"), "test.pptx"),
                },
                content_type="multipart/form-data",
            )
            self.assertEqual(response.status_code, 400)
            remove_file.assert_called_once_with(
                str(Path("test-data") / "custom_catalogue_uploads" / "_incoming" / "test.pptx")
            )


if __name__ == "__main__":
    unittest.main()
