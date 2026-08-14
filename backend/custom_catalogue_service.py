import json
import math
import re
import shutil
import threading
from datetime import datetime
from io import BytesIO
from pathlib import Path

import fitz
from PIL import Image, ImageFilter, ImageOps
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml import parse_xml


BASE_DIR = Path(__file__).resolve().parent.parent
DATA_DIR = BASE_DIR / "data"
FRONTEND_DIR = BASE_DIR / "frontend"
CUSTOM_STATIC_DIR = FRONTEND_DIR / "static" / "custom_catalogues"
UPLOAD_DIR = DATA_DIR / "custom_catalogue_uploads"
REGISTRY_PATH = DATA_DIR / "custom_catalogues.json"
_CUSTOM_ROOM_LOCK = threading.RLock()

NS = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
SCHEME_COLORS = {"accent4": "8064A2"}
RISK_BY_PREFIX = {"M": "medium", "H": "high", "L": "low", "O": "office"}
EXPLICIT_ROOM_CODE = re.compile(r"^([A-Za-z]{1,4})[-\s]?(\d{1,3})\s*[_:|-]\s*(.+)$")
DEFAULT_COLORS = {
    "medium": "0EA5E9",
    "high": "22C55E",
    "low": "EAB308",
    "office": "8064A2",
    "custom": "2563EB",
}


def rounded(value):
    return round(value, 4)


def slugify(value):
    slug = re.sub(r"[^a-z0-9]+", "-", value.strip().lower())
    slug = re.sub(r"-+", "-", slug).strip("-")
    return slug or "custom-catalogue"


def unique_slug(title):
    base_slug = slugify(title)
    registry = load_custom_catalogues()
    existing = {item.get("slug") for item in registry}
    slug = base_slug
    counter = 2
    while slug in existing:
        slug = f"{base_slug}-{counter}"
        counter += 1
    return slug


def load_custom_catalogues():
    if not REGISTRY_PATH.exists():
        return []
    try:
        with REGISTRY_PATH.open("r", encoding="utf-8") as file:
            payload = json.load(file)
    except (json.JSONDecodeError, OSError):
        return []
    return payload if isinstance(payload, list) else []


def save_custom_catalogues(items):
    DATA_DIR.mkdir(parents=True, exist_ok=True)
    temp_path = REGISTRY_PATH.with_suffix(".json.tmp")
    with temp_path.open("w", encoding="utf-8") as file:
        json.dump(items, file, indent=2)
    temp_path.replace(REGISTRY_PATH)


def get_custom_catalogue(slug):
    normalized_slug = slugify(slug)
    return next(
        (item for item in load_custom_catalogues() if item.get("slug") == normalized_slug),
        None,
    )


def parse_room_shape(shape_name, fallback_code=None):
    shape_name = str(shape_name or "").strip()
    explicit_match = EXPLICIT_ROOM_CODE.match(shape_name)
    if explicit_match:
        prefix, number, raw_name = explicit_match.groups()
        code = f"{prefix.upper()}-{int(number):02d}"
    elif fallback_code:
        code = fallback_code
        raw_name = shape_name
    else:
        return None

    room_name = raw_name.replace("_", " ").replace(" . ", " ").strip()
    room_name = " ".join(room_name.split())
    if not room_name:
        return None

    risk = RISK_BY_PREFIX.get(code[:1], "custom")
    return code, room_name, risk


def find_shape(slide, shape_name):
    return next((shape for shape in slide.shapes if shape.name == shape_name), None)


def find_largest_picture(slide):
    pictures = [
        shape
        for shape in slide.shapes
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
    ]
    if not pictures:
        return None
    return max(pictures, key=lambda shape: shape.width * shape.height)


def _intersection_area(first_bounds, second_bounds):
    left = max(first_bounds[0], second_bounds[0])
    top = max(first_bounds[1], second_bounds[1])
    right = min(first_bounds[2], second_bounds[2])
    bottom = min(first_bounds[3], second_bounds[3])
    return max(0, right - left) * max(0, bottom - top)


def find_map_pictures(slide, picture_name=""):
    pictures = [
        shape
        for shape in slide.shapes
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
    ]
    if not pictures:
        return []

    if picture_name:
        requested = find_shape(slide, picture_name)
        if requested is not None and requested.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return [requested]

    room_shapes = [
        shape
        for shape in slide.shapes
        if shape.shape_type in {MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.FREEFORM}
        and parse_room_shape(shape.name, "C-01")
    ]
    if room_shapes:
        room_bounds = [get_rendered_shape_bounds(shape) for shape in room_shapes]
        overlapping_pictures = [
            picture
            for picture in pictures
            if any(
                _intersection_area(get_rendered_shape_bounds(picture), bounds) > 0
                for bounds in room_bounds
            )
        ]
        if overlapping_pictures:
            return overlapping_pictures

    return [max(pictures, key=lambda shape: shape.width * shape.height)]


def crop_picture_image(picture):
    image = Image.open(BytesIO(picture.image.blob)).convert("RGBA")
    width, height = image.size
    box = (
        round(picture.crop_left * width),
        round(picture.crop_top * height),
        round((1 - picture.crop_right) * width),
        round((1 - picture.crop_bottom) * height),
    )
    return image.crop(box)


def get_shape_transform(shape):
    transform = shape.element.find(".//a:xfrm", namespaces=NS)
    if transform is None:
        return 0, False, False
    rotation = int(transform.get("rot") or 0) / 60000
    flip_h = str(transform.get("flipH") or "").lower() in {"1", "true"}
    flip_v = str(transform.get("flipV") or "").lower() in {"1", "true"}
    return rotation, flip_h, flip_v


def transform_shape_point(shape, local_x, local_y):
    rotation, flip_h, flip_v = get_shape_transform(shape)
    if flip_h:
        local_x = 1 - local_x
    if flip_v:
        local_y = 1 - local_y

    center_x = shape.left + shape.width / 2
    center_y = shape.top + shape.height / 2
    offset_x = (local_x - 0.5) * shape.width
    offset_y = (local_y - 0.5) * shape.height
    radians = math.radians(rotation)
    return (
        center_x + offset_x * math.cos(radians) - offset_y * math.sin(radians),
        center_y + offset_x * math.sin(radians) + offset_y * math.cos(radians),
    )


def get_rendered_shape_bounds(shape):
    corners = [
        transform_shape_point(shape, local_x, local_y)
        for local_x, local_y in ((0, 0), (1, 0), (1, 1), (0, 1))
    ]
    xs = [point[0] for point in corners]
    ys = [point[1] for point in corners]
    return min(xs), min(ys), max(xs), max(ys)


def normalize_map_pictures(picture_or_pictures):
    if isinstance(picture_or_pictures, (list, tuple)):
        return [picture for picture in picture_or_pictures if picture is not None]
    return [picture_or_pictures] if picture_or_pictures is not None else []


def get_map_bounds(picture_or_pictures):
    pictures = normalize_map_pictures(picture_or_pictures)
    if not pictures:
        raise ValueError("A map picture is required.")
    bounds = [get_rendered_shape_bounds(picture) for picture in pictures]
    return (
        min(item[0] for item in bounds),
        min(item[1] for item in bounds),
        max(item[2] for item in bounds),
        max(item[3] for item in bounds),
    )


def export_floorplan(picture_or_pictures, output_path, target_width):
    pictures = normalize_map_pictures(picture_or_pictures)
    left, top, right, bottom = get_map_bounds(pictures)
    rendered_width = max(1, right - left)
    rendered_height = max(1, bottom - top)
    scale = target_width / rendered_width
    target_height = max(1, round(target_width * rendered_height / rendered_width))
    canvas = Image.new("RGBA", (target_width, target_height), "white")

    for picture in pictures:
        image = crop_picture_image(picture)
        rotation, flip_h, flip_v = get_shape_transform(picture)
        picture_left, picture_top, picture_right, picture_bottom = get_rendered_shape_bounds(picture)
        unrotated_size = (
            max(1, round(picture.width * scale)),
            max(1, round(picture.height * scale)),
        )
        image = image.resize(unrotated_size, Image.Resampling.LANCZOS)
        if flip_h:
            image = ImageOps.mirror(image)
        if flip_v:
            image = ImageOps.flip(image)
        if rotation % 360:
            image = image.rotate(-rotation, expand=True, resample=Image.Resampling.BICUBIC)

        expected_size = (
            max(1, round((picture_right - picture_left) * scale)),
            max(1, round((picture_bottom - picture_top) * scale)),
        )
        if image.size != expected_size:
            image = image.resize(expected_size, Image.Resampling.LANCZOS)
        canvas.alpha_composite(image, (
            round((picture_left - left) * scale),
            round((picture_top - top) * scale),
        ))

    image = canvas.convert("RGB").filter(ImageFilter.UnsharpMask(radius=1.2, percent=135, threshold=2))
    output_path.parent.mkdir(parents=True, exist_ok=True)
    image.save(output_path, "PNG", optimize=True)
    return target_width, target_height


def shape_point_to_picture(point, path_width, path_height, shape, picture_or_pictures):
    local_x = int(point.get("x")) / path_width
    local_y = int(point.get("y")) / path_height
    slide_x, slide_y = transform_shape_point(shape, local_x, local_y)
    picture_left, picture_top, picture_right, picture_bottom = get_map_bounds(picture_or_pictures)
    return (
        rounded((slide_x - picture_left) / (picture_right - picture_left) * 100),
        rounded((slide_y - picture_top) / (picture_bottom - picture_top) * 100),
    )


def get_shape_rotation(shape):
    return get_shape_transform(shape)[0]


def extract_rotated_rect_path(shape, picture_or_pictures):
    rotation = get_shape_rotation(shape)
    if rotation % 360 == 0:
        return None

    picture_left, picture_top, picture_right, picture_bottom = get_map_bounds(picture_or_pictures)
    picture_width = picture_right - picture_left
    picture_height = picture_bottom - picture_top
    corners = [
        (
            rounded((slide_x - picture_left) / picture_width * 100),
            rounded((slide_y - picture_top) / picture_height * 100),
        )
        for slide_x, slide_y in (
            transform_shape_point(shape, local_x, local_y)
            for local_x, local_y in ((0, 0), (1, 0), (1, 1), (0, 1))
        )
    ]

    return " ".join(
        [f"M {corners[0][0]} {corners[0][1]}"]
        + [f"L {x} {y}" for x, y in corners[1:]]
        + ["Z"]
    )


def extract_svg_path(shape, picture_or_pictures):
    path = shape.element.find(".//a:custGeom/a:pathLst/a:path", namespaces=NS)
    if path is None:
        return extract_rotated_rect_path(shape, picture_or_pictures)

    path_width = int(path.get("w") or shape.width)
    path_height = int(path.get("h") or shape.height)
    commands = []

    for command in path:
        tag = command.tag.rsplit("}", 1)[-1]
        points = command.findall("a:pt", namespaces=NS)
        coords = [
            shape_point_to_picture(point, path_width, path_height, shape, picture_or_pictures)
            for point in points
        ]

        if tag == "moveTo" and coords:
            commands.append(f"M {coords[0][0]} {coords[0][1]}")
        elif tag == "lnTo" and coords:
            commands.append(f"L {coords[0][0]} {coords[0][1]}")
        elif tag == "cubicBezTo" and len(coords) == 3:
            commands.append("C " + " ".join(f"{x} {y}" for x, y in coords))
        elif tag == "quadBezTo" and len(coords) == 2:
            commands.append("Q " + " ".join(f"{x} {y}" for x, y in coords))
        elif tag == "close":
            commands.append("Z")

    return " ".join(commands) or None


def get_theme_colors(slide):
    colors = dict(SCHEME_COLORS)
    try:
        master = slide.slide_layout.slide_master
        theme_part = master.part.part_related_by(RT.THEME)
        theme = parse_xml(theme_part.blob)
        color_scheme = theme.find(".//a:themeElements/a:clrScheme", namespaces=NS)
        if color_scheme is not None:
            for scheme_entry in color_scheme:
                if len(scheme_entry) == 0:
                    continue
                color_node = scheme_entry[0]
                tag = color_node.tag.rsplit("}", 1)[-1]
                value = color_node.get("val") if tag == "srgbClr" else color_node.get("lastClr")
                if value:
                    colors[scheme_entry.tag.rsplit("}", 1)[-1]] = value.upper()

        color_map = master.element.find(".//p:clrMap", namespaces={
            **NS,
            "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
        })
        if color_map is not None:
            for alias, target in color_map.attrib.items():
                if target in colors:
                    colors[alias] = colors[target]
    except (AttributeError, KeyError, ValueError):
        pass
    return colors


def resolve_color(container, default, theme_colors=None):
    if container is None or len(container) == 0:
        return default
    color = container[0]
    tag = color.tag.rsplit("}", 1)[-1]
    if tag == "srgbClr":
        return color.get("val") or default
    if tag == "schemeClr":
        return (theme_colors or SCHEME_COLORS).get(color.get("val"), default)
    if tag == "sysClr":
        return color.get("lastClr") or default
    return default


def extract_style(shape, risk, theme_colors=None):
    shape_properties = shape.element.find("p:spPr", namespaces={
        **NS,
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    })
    solid_fill = shape_properties.find("a:solidFill", namespaces=NS) if shape_properties is not None else None
    line = shape_properties.find("a:ln", namespaces=NS) if shape_properties is not None else None
    line_fill = line.find("a:solidFill", namespaces=NS) if line is not None else None
    fill_color = resolve_color(
        solid_fill,
        DEFAULT_COLORS.get(risk, DEFAULT_COLORS["custom"]),
        theme_colors,
    )
    stroke_color = resolve_color(line_fill, fill_color, theme_colors)
    color_node = solid_fill[0] if solid_fill is not None and len(solid_fill) else None
    alpha = color_node.find("a:alpha", namespaces=NS) if color_node is not None else None
    fill_opacity = int(alpha.get("val")) / 100000 if alpha is not None else 1
    line_width = int(line.get("w") or 12700) / 12700 if line is not None else 1
    return {
        "fillColor": f"#{fill_color}",
        "fillOpacity": rounded(fill_opacity),
        "strokeColor": f"#{stroke_color}",
        "strokeWidth": rounded(line_width),
    }


def extract_room(shape, picture_or_pictures, fallback_code=None, theme_colors=None):
    parsed = parse_room_shape(shape.name, fallback_code)
    if not parsed:
        return None

    code, room_name, risk = parsed
    shape_left, shape_top, shape_right, shape_bottom = get_rendered_shape_bounds(shape)
    picture_left, picture_top, picture_right, picture_bottom = get_map_bounds(picture_or_pictures)
    picture_width = picture_right - picture_left
    picture_height = picture_bottom - picture_top
    return {
        "code": code,
        "name": room_name,
        "risk": risk,
        "interactive": True,
        "left": rounded((shape_left - picture_left) / picture_width * 100),
        "top": rounded((shape_top - picture_top) / picture_height * 100),
        "width": rounded((shape_right - shape_left) / picture_width * 100),
        "height": rounded((shape_bottom - shape_top) / picture_height * 100),
        "svgPath": extract_svg_path(shape, picture_or_pictures),
        **extract_style(shape, risk, theme_colors),
    }


def extract_rooms_from_slide(slide, picture_or_pictures, start_index=1):
    rooms = {}
    generated_index = max(1, int(start_index or 1))
    theme_colors = get_theme_colors(slide)
    for shape in slide.shapes:
        if shape.shape_type not in {
            MSO_SHAPE_TYPE.AUTO_SHAPE,
            MSO_SHAPE_TYPE.FREEFORM,
        }:
            continue
        fallback_code = f"C-{generated_index:02d}"
        room = extract_room(shape, picture_or_pictures, fallback_code, theme_colors)
        if room:
            generated_index += 1
            if room["code"] in rooms:
                raise ValueError(
                    f"Duplicate room code {room['code']} in PowerPoint shape names. "
                    "Give each coded shape a unique room code."
                )
            rooms[room["code"]] = room
    return dict(sorted(rooms.items()))


def parse_slide_numbers(value, slide_count, default=1):
    if isinstance(value, (list, tuple, set)):
        requested = list(value)
    else:
        text = str(value or default).strip().lower()
        if text in {"all", "*"}:
            return list(range(1, slide_count + 1))
        requested = []
        for token in re.split(r"[\s,;]+", text):
            if not token:
                continue
            if "-" in token:
                parts = token.split("-", 1)
                try:
                    start, end = int(parts[0]), int(parts[1])
                except ValueError as error:
                    raise ValueError("Slides must be 'all', comma-separated numbers, or ranges such as 1-3.") from error
                if end < start:
                    start, end = end, start
                requested.extend(range(start, end + 1))
            else:
                try:
                    requested.append(int(token))
                except ValueError as error:
                    raise ValueError("Slides must be 'all', comma-separated numbers, or ranges such as 1-3.") from error

    slide_numbers = []
    for raw_number in requested:
        try:
            number = int(raw_number)
        except (TypeError, ValueError) as error:
            raise ValueError("Every slide number must be a whole number.") from error
        if number < 1 or number > slide_count:
            raise ValueError(f"The PPT only has {slide_count} slides; slide {number} is unavailable.")
        if number not in slide_numbers:
            slide_numbers.append(number)
    if not slide_numbers:
        raise ValueError("Choose at least one PowerPoint slide.")
    return slide_numbers


def get_pdf_page_object_order(pdf_bytes):
    page_objects = []

    for match in re.finditer(rb"/Type /Page\b", pdf_bytes):
        object_marker = pdf_bytes.rfind(b" obj", 0, match.start())
        if object_marker == -1:
            continue

        object_start = pdf_bytes.rfind(b"\n", 0, object_marker)
        object_header = pdf_bytes[object_start + 1:object_marker + 4]
        object_match = re.search(rb"(\d+)\s+0\s+obj", object_header)
        if object_match:
            page_objects.append(int(object_match.group(1)))

    return {object_id: index + 1 for index, object_id in enumerate(page_objects)}


def get_pdf_named_destinations(pdf_bytes, page_by_object):
    named_destinations = {}
    destination_marker = pdf_bytes.rfind(b"/h.")
    if destination_marker == -1:
        return named_destinations

    object_start = pdf_bytes.rfind(b" obj", 0, destination_marker)
    if object_start == -1:
        return named_destinations

    object_start = pdf_bytes.rfind(b"\n", 0, object_start)
    object_end = pdf_bytes.find(b"endobj", destination_marker)
    if object_end == -1:
        return named_destinations

    object_body = pdf_bytes[object_start:object_end]
    for destination_match in re.finditer(
        rb"/(h\.[^\s\[]+)\s*\[(\d+)\s+0\s+R",
        object_body,
    ):
        destination_name = destination_match.group(1).decode("latin-1")
        page_object = int(destination_match.group(2))
        page_number = page_by_object.get(page_object)
        if page_number:
            named_destinations[destination_name] = page_number

    return named_destinations


def get_pdf_contents_links(pdf_bytes):
    links = []
    with fitz.open(stream=pdf_bytes, filetype="pdf") as document:
        for source_page in range(min(document.page_count, 10)):
            page = document[source_page]
            for link in page.get_links():
                rectangle = link.get("from")
                destination_page = link.get("page")
                if (
                    not rectangle
                    or not isinstance(destination_page, int)
                    or destination_page <= source_page
                    or rectangle.width < 80
                ):
                    continue
                text = normalize_pdf_link_text(page.get_textbox(rectangle))
                if text:
                    links.append({
                        "page": destination_page + 1,
                        "text": text,
                        "left": rectangle.x0,
                        "width": rectangle.width,
                    })
    return links


def normalize_pdf_link_text(value):
    text = (value or "").replace("\u200b", " ").replace("\ufeff", " ").lower()
    return " ".join(re.sub(r"[^a-z0-9]+", " ", text).split())


def match_pdf_pages_to_rooms(rooms, links):
    updated_rooms = {code: dict(room) for code, room in rooms.items()}
    available_links = list(enumerate(links))
    used_links = set()

    room_entries = sorted(
        updated_rooms.items(),
        key=lambda item: len(normalize_pdf_link_text(item[1].get("name"))),
        reverse=True,
    )
    for code, room in room_entries:
        code_text = normalize_pdf_link_text(code)
        name_text = normalize_pdf_link_text(room.get("name"))
        best_match = None
        best_score = 0
        for link_index, link in available_links:
            if link_index in used_links:
                continue
            link_text = link["text"]
            score = 0
            if code_text and code_text in link_text:
                score = 1000 + len(code_text)
            if name_text and (name_text in link_text or link_text in name_text):
                score = max(score, 500 + min(len(name_text), len(link_text)))
            if score > best_score:
                best_match = (link_index, link)
                best_score = score
        if best_match:
            link_index, link = best_match
            used_links.add(link_index)
            room["page"] = link["page"]

    return updated_rooms, len(used_links)


def match_pdf_headers_to_rooms(rooms, pdf_bytes, max_header_lines=3):
    """Match room names against the visible heading at the top of PDF pages.

    Only the first few non-empty lines are considered. This deliberately avoids
    matching the room names repeated later on a contents page.
    """
    updated_rooms = {code: dict(room) for code, room in rooms.items()}
    matched_codes = set()
    with fitz.open(stream=pdf_bytes, filetype="pdf") as document:
        page_headers = []
        for page_index, page in enumerate(document):
            lines = [
                normalize_pdf_link_text(line)
                for line in page.get_text().splitlines()
                if normalize_pdf_link_text(line)
            ][:max_header_lines]
            page_headers.append((page_index + 1, lines))

    room_entries = sorted(
        updated_rooms.items(),
        key=lambda item: len(normalize_pdf_link_text(item[1].get("name"))),
        reverse=True,
    )
    for code, room in room_entries:
        code_text = normalize_pdf_link_text(code)
        name_text = normalize_pdf_link_text(room.get("name"))
        candidates = []
        for page_number, header_lines in page_headers:
            best_score = 0
            for line_index, header_text in enumerate(header_lines):
                score = 0
                if name_text and (name_text in header_text or header_text in name_text):
                    score = 1000 + min(len(name_text), len(header_text)) - (line_index * 10)
                if code_text and code_text in header_text:
                    score = max(score, 2000 + len(code_text) - (line_index * 10))
                best_score = max(best_score, score)
            if best_score:
                candidates.append((best_score, page_number))
        if candidates:
            candidates.sort(reverse=True)
            room["page"] = candidates[0][1]
            matched_codes.add(code)

    return updated_rooms, len(matched_codes)


def get_pdf_contents_link_pages(pdf_bytes):
    try:
        link_pages = [
            link["page"]
            for link in get_pdf_contents_links(pdf_bytes)
            if abs(link["left"] - 108) < 0.5 and link["width"] > 250
        ]
        if link_pages:
            return link_pages
    except Exception:
        # Keep support for unusual PDFs by falling back to the original
        # byte-level named-destination parser below.
        pass

    page_by_object = get_pdf_page_object_order(pdf_bytes)
    named_destinations = get_pdf_named_destinations(pdf_bytes, page_by_object)
    link_pages = []

    for object_number in range(1, 80):
        object_marker = f"\n{object_number} 0 obj".encode("ascii")
        object_start = pdf_bytes.find(object_marker)
        if object_start == -1 and pdf_bytes.startswith(object_marker[1:]):
            object_start = 0
        if object_start == -1:
            continue

        object_end = pdf_bytes.find(b"endobj", object_start)
        if object_end == -1:
            continue

        object_body = pdf_bytes[object_start:object_end]
        rect_match = re.search(rb"/Rect \[([^\]]+)\]", object_body)
        destination_match = re.search(rb"/Dest /(h\.[^\s>/]+)", object_body)
        if not rect_match or not destination_match:
            continue

        try:
            rect_values = [float(value) for value in rect_match.group(1).split()]
        except ValueError:
            continue

        destination_name = destination_match.group(1).decode("latin-1")
        page_number = named_destinations.get(destination_name)
        if (
            len(rect_values) >= 4
            and abs(rect_values[0] - 108) < 0.5
            and (rect_values[2] - rect_values[0]) > 250
            and page_number
        ):
            link_pages.append(page_number)

    return link_pages


def apply_pdf_pages_to_rooms(rooms, pdf_bytes):
    if not pdf_bytes:
        return rooms

    updated_rooms = {code: dict(room) for code, room in rooms.items()}
    match_count = 0
    try:
        updated_rooms, match_count = match_pdf_pages_to_rooms(
            rooms,
            get_pdf_contents_links(pdf_bytes),
        )
    except Exception:
        pass

    # A page title is the most reliable source of truth. It validates and, when
    # necessary, corrects destinations inferred from the Google Docs contents.
    try:
        header_rooms, header_match_count = match_pdf_headers_to_rooms(updated_rooms, pdf_bytes)
        if header_match_count:
            return header_rooms
    except Exception:
        pass

    if match_count:
        return updated_rooms

    toc_pages = get_pdf_contents_link_pages(pdf_bytes)
    if len(toc_pages) < len(rooms):
        return rooms

    updated_rooms = {
        code: dict(room)
        for code, room in rooms.items()
    }
    for index, code in enumerate(sorted(updated_rooms)):
        updated_rooms[code]["page"] = toc_pages[index]
    return updated_rooms


def apply_custom_room_overrides(rooms, overrides):
    updated_rooms = {code: dict(room) for code, room in rooms.items()}
    for code, room in updated_rooms.items():
        room.setdefault("defaultName", room.get("name", ""))
        room.setdefault("defaultPage", room.get("page"))
        override = (overrides or {}).get(code, {})
        if override.get("name"):
            room["name"] = override["name"]
        if override.get("page") not in ("", None):
            try:
                room["page"] = max(1, int(override["page"]))
            except (TypeError, ValueError):
                pass
    return updated_rooms


def create_custom_catalogue(title, doc_url, pptx_path, pdf_bytes, options=None):
    options = options or {}
    slug = slugify(options.get("slug")) if options.get("slug") else unique_slug(title)
    target_width = int(options.get("target_width") or 3600)
    picture_name = (options.get("picture_name") or "").strip()

    if target_width < 800 or target_width > 10000:
        raise ValueError("Target width must be between 800 and 10000 pixels.")

    presentation = Presentation(str(pptx_path))
    slide_value = options.get("slide_numbers", options.get("slide_number", 1))
    slide_numbers = parse_slide_numbers(slide_value, len(presentation.slides))
    scan_all_slides = str(slide_value or "").strip().lower() in {"all", "*"}

    extracted_layouts = []
    rooms = {}
    next_generated_index = 1
    for slide_number in slide_numbers:
        slide = presentation.slides[slide_number - 1]
        pictures = find_map_pictures(slide, picture_name)
        if not pictures:
            if scan_all_slides:
                continue
            raise ValueError(f"Could not find the map picture on slide {slide_number}.")

        slide_rooms = extract_rooms_from_slide(slide, pictures, next_generated_index)
        if not slide_rooms:
            if scan_all_slides:
                continue
            raise ValueError(
                f"No room shapes were found on slide {slide_number}. Add named AutoShape or Freeform overlays "
                "on top of the map picture. Friendly names such as Kitchen Area are supported."
            )

        layout_key = f"slide-{slide_number}"
        for code, room in slide_rooms.items():
            if code in rooms:
                raise ValueError(
                    f"Duplicate room code {code} across the selected slides. "
                    "Give every room shape a unique code or friendly name."
                )
            rooms[code] = {
                **room,
                "layoutKey": layout_key,
                "slideNumber": slide_number,
            }
        next_generated_index += len(slide_rooms)
        extracted_layouts.append({
            "key": layout_key,
            "label": f"Map {len(extracted_layouts) + 1} · Slide {slide_number}",
            "slide_number": slide_number,
            "picture_name": " + ".join(picture.name for picture in pictures),
            "picture_names": [picture.name for picture in pictures],
            "pictures": pictures,
            "room_codes": list(slide_rooms),
        })

    if not extracted_layouts:
        raise ValueError(
            "No usable map slides were found. Each map slide needs a picture plus at least one named "
            "AutoShape or Freeform overlay."
        )

    rooms = dict(sorted(rooms.items()))
    rooms = apply_pdf_pages_to_rooms(rooms, pdf_bytes)
    room_overrides = options.get("room_overrides") or {}
    rooms = apply_custom_room_overrides(rooms, room_overrides)

    catalogue_dir = CUSTOM_STATIC_DIR / slug
    catalogue_dir.mkdir(parents=True, exist_ok=True)

    saved_pptx_path = UPLOAD_DIR / f"{slug}.pptx"
    UPLOAD_DIR.mkdir(parents=True, exist_ok=True)
    if Path(pptx_path).resolve() != saved_pptx_path.resolve():
        shutil.copyfile(pptx_path, saved_pptx_path)

    layouts = []
    for index, extracted_layout in enumerate(extracted_layouts):
        image_name = "layout.png" if index == 0 else f"layout-slide-{extracted_layout['slide_number']}.png"
        image_width, image_height = export_floorplan(
            extracted_layout["pictures"],
            catalogue_dir / image_name,
            target_width,
        )
        layouts.append({
            "key": extracted_layout["key"],
            "label": extracted_layout["label"],
            "slide_number": extracted_layout["slide_number"],
            "picture_name": extracted_layout["picture_name"],
            "picture_names": extracted_layout["picture_names"],
            "image": f"/static/custom_catalogues/{slug}/{image_name}",
            "width": image_width,
            "height": image_height,
            "aspect": rounded(image_width / image_height),
            "room_codes": extracted_layout["room_codes"],
        })

    if pdf_bytes:
        with (catalogue_dir / "catalogue.pdf").open("wb") as file:
            file.write(pdf_bytes)

    generated_at = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    catalogue = {
        "slug": slug,
        "title": title.strip(),
        "doc_url": doc_url.strip(),
        "created_at": options.get("created_at") or generated_at,
        "updated_at": generated_at,
        "last_checked": options.get("last_checked", ""),
        "last_updated": options.get("last_updated") or generated_at,
        "refresh_error": options.get("refresh_error", ""),
        "auto_refresh": bool(doc_url.strip()),
        "refresh_interval_minutes": options.get("refresh_interval_minutes", 5),
        "source_pptx": str(saved_pptx_path.relative_to(BASE_DIR)).replace("\\", "/"),
        "slide_number": layouts[0]["slide_number"],
        "slide_numbers": [layout["slide_number"] for layout in layouts],
        "picture_name": layouts[0]["picture_name"],
        "picture_names": [layout["picture_name"] for layout in layouts],
        "room_count": len(rooms),
        "room_overrides": room_overrides,
        "layout": {
            "image": layouts[0]["image"],
            "pdf": f"/static/custom_catalogues/{slug}/catalogue.pdf" if pdf_bytes else "",
            "width": layouts[0]["width"],
            "height": layouts[0]["height"],
            "aspect": layouts[0]["aspect"],
        },
        "layouts": layouts,
        "rooms": rooms,
    }

    with (catalogue_dir / "catalogue.json").open("w", encoding="utf-8") as file:
        json.dump(catalogue, file, indent=2)

    registry = [item for item in load_custom_catalogues() if item.get("slug") != slug]
    registry.append(catalogue)
    registry.sort(key=lambda item: item.get("created_at", ""), reverse=True)
    save_custom_catalogues(registry)
    return catalogue


def rebuild_custom_catalogue(slug, pdf_bytes=None):
    current = get_custom_catalogue(slug)
    if not current:
        raise ValueError("Custom catalogue not found.")

    source_pptx = BASE_DIR / current.get("source_pptx", "")
    if not source_pptx.is_file():
        raise ValueError("The saved PowerPoint source for this catalogue is missing.")

    if pdf_bytes is None:
        pdf_path = CUSTOM_STATIC_DIR / current["slug"] / "catalogue.pdf"
        pdf_bytes = pdf_path.read_bytes() if pdf_path.is_file() else b""
    return create_custom_catalogue(
        current.get("title") or current["slug"],
        current.get("doc_url", ""),
        source_pptx,
        pdf_bytes,
        {
            "slug": current["slug"],
            "created_at": current.get("created_at"),
            "slide_numbers": current.get("slide_numbers", [current.get("slide_number", 1)]),
            "picture_name": current.get("picture_name", ""),
            "target_width": (current.get("layout") or {}).get("width", 3600),
            "last_checked": current.get("last_checked", ""),
            "last_updated": current.get("last_updated", ""),
            "refresh_error": current.get("refresh_error", ""),
            "refresh_interval_minutes": current.get("refresh_interval_minutes", 5),
            "room_overrides": current.get("room_overrides", {}),
        },
    )


def upsert_custom_catalogue_room(slug, data):
    normalized_slug = slugify(slug)
    room_code = str((data or {}).get("code") or "").strip().upper()
    if not room_code:
        raise ValueError("Room code is required")

    with _CUSTOM_ROOM_LOCK:
        catalogues = load_custom_catalogues()
        catalogue = next((item for item in catalogues if item.get("slug") == normalized_slug), None)
        if not catalogue:
            raise ValueError("Catalogue not found")
        room = (catalogue.get("rooms") or {}).get(room_code)
        if not room:
            raise ValueError("Room not found in this catalogue")

        room.setdefault("defaultName", room.get("name", ""))
        room.setdefault("defaultPage", room.get("page"))
        overrides = catalogue.setdefault("room_overrides", {})
        override = overrides.setdefault(room_code, {})
        name = str((data or {}).get("name") or "").strip()
        page = (data or {}).get("page")

        if name and name != room.get("defaultName"):
            override["name"] = name
            room["name"] = name
        else:
            override.pop("name", None)
            room["name"] = room.get("defaultName", "")

        if page not in ("", None):
            try:
                page_number = max(1, int(page))
            except (TypeError, ValueError) as error:
                raise ValueError("Catalogue page must be a positive number") from error
            if page_number != room.get("defaultPage"):
                override["page"] = page_number
            else:
                override.pop("page", None)
            room["page"] = page_number
        else:
            override.pop("page", None)
            room["page"] = room.get("defaultPage")

        override["updated_at"] = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        if not any(key in override for key in ("name", "page")):
            overrides.pop(room_code, None)
        catalogue["updated_at"] = datetime.now().strftime("%Y-%m-%d %H:%M:%S")

        catalogue_dir = CUSTOM_STATIC_DIR / normalized_slug
        catalogue_dir.mkdir(parents=True, exist_ok=True)
        with (catalogue_dir / "catalogue.json").open("w", encoding="utf-8") as file:
            json.dump(catalogue, file, indent=2)
        save_custom_catalogues(catalogues)
        return {
            "catalogue_id": f"custom:{normalized_slug}",
            "code": room_code,
            "name": room.get("name", ""),
            "default_name": room.get("defaultName", ""),
            "page": room.get("page"),
        }
